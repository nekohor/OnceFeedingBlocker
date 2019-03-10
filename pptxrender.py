from pptx import Presentation
from pptx.util import Inches, Pt


class PPTXRender:

    def __init__(self, ctx, rec):
        self.ctx = ctx
        self.rec = rec

    def build_pptx(self):
        prs = Presentation()

        for tag in self.ctx.tags:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 每页标题
            left = top = Inches(0.1)
            width = Inches(2)
            height = Inches(1)  # 预设位置及大小
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = "{}{}{}".format(
                self.rec.get_item_name(),
                self.ctx.lang_map[tag],
                self.ctx.cn_name
            )

            # 图片罗列
            cfg_table = self.ctx.plot_config
            for idx in cfg_table.index:
                cfg = cfg_table.loc[idx]

                plot_title = cfg["title"]
                plot_file_name = (
                    self.rec.get_item_file_name(
                        "{}_{}.png".format(tag, plot_title)
                    )
                )

                left = Inches(cfg["left"])
                top = Inches(cfg["top"])
                width = Inches(cfg["width"])
                height = Inches(cfg["height"])

                pic = slide.shapes.add_picture(
                    plot_file_name,
                    left, top, width, height)
                pic

            # 原因表格绘制
            rows, cols, = 3, 8
            left, top = Inches(0.1), Inches(5.8)
            width, height = Inches(9.5), Inches(1.6)
            table = slide.shapes.add_table(
                rows, cols, left, top, width, height).table
            table.columns[0].width = Inches(0.8)
            table.columns[1].width = Inches(1.5)
            table.columns[2].width = Inches(1.5)
            table.columns[3].width = Inches(1.5)
            table.columns[4].width = Inches(1.5)
            table.columns[5].width = Inches(1)
            table.columns[6].width = Inches(0.8)
            table.columns[7].width = Inches(1.2)

            table.cell(0, 0).text = '序号'
            table.cell(1, 0).text = '1'
            table.cell(2, 0).text = '2'

            table.cell(0, 1).text = '主要原因'
            table.cell(0, 2).text = '对策'
            table.cell(0, 3).text = '目标'
            table.cell(0, 4).text = '措施'
            table.cell(0, 5).text = '负责人'
            table.cell(0, 6).text = '地点'
            table.cell(0, 7).text = '完成时间'

            prs.save(self.rec.get_item_file_name('四版图.pptx'))
